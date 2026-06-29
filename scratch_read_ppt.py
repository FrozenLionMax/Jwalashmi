import collections
import collections.abc
from pptx import Presentation
import sys
import json

def inspect_ppt(filepath):
    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        return

    output = []
    for i, slide in enumerate(prs.slides):
        slide_info = {"slide_number": i + 1, "shapes": []}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.replace('\n', ' ').replace('\r', '')
            if text.strip():
                slide_info["shapes"].append({"shape_name": shape.name, "shape_id": shape.shape_id, "text": text[:50]})
        output.append(slide_info)

    print(json.dumps(output, indent=2))

if __name__ == "__main__":
    inspect_ppt("[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx")
